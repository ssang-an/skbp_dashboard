from __future__ import annotations

import json
import os
from collections import Counter
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "presentation_assets"
OUT = ROOT / "SKBP_Pipeline_Dashboard_전사발표_7분_이미지생성버전.pptx"
SCRIPT_OUT = ROOT / "SKBP_Pipeline_Dashboard_7분_발표원고.md"
ASSET_DIR.mkdir(exist_ok=True)

# ---------- Data ----------
with (ROOT / "json" / "pipeline-records.json").open(encoding="utf-8") as f:
    records = json.load(f)

def n(x):
    try:
        return float(x)
    except Exception:
        return None

criteria_keys = [
    ("target_relevance", "TR"),
    ("competitive_landscape", "Comp"),
    ("moa_validity", "MOA"),
    ("platform_attractiveness", "Plat"),
    ("expansion_potential", "Exp"),
    ("data_maturity", "Data"),
    ("marketability", "Market"),
]

themes, clusters, countries, statuses, stages, target_scores, totals = Counter(), Counter(), Counter(), Counter(), Counter(), Counter(), []
for r in records:
    js = r.get("json_summary") or {}
    st = r.get("structured_table") or {}
    sc = r.get("scoring") or {}
    hf = r.get("hard_filter") or {}
    themes[js.get("theme") or "Unknown"] += 1
    clusters[js.get("cluster") or "Unknown"] += 1
    countries[js.get("company_country") or st.get("company_country") or "Unknown"] += 1
    statuses[hf.get("status") or hf.get("overall_result") or "REVIEW"] += 1
    stages[st.get("development_stage") or "Unknown"] += 1
    tr = n(js.get("target_relevance_score") or (((sc.get("criteria") or {}).get("target_relevance") or {}).get("score")))
    target_scores[str(int(tr)) if tr is not None else "미평가"] += 1
    total = n(sc.get("total_score"))
    if total is not None:
        totals.append(total)

pass_select = statuses.get("SELECT", 0) + statuses.get("PASS", 0)
avg_total = sum(totals) / len(totals) if totals else 0
avg_target = sum(float(k) * v for k, v in target_scores.items() if k.isdigit()) / sum(v for k, v in target_scores.items() if k.isdigit())

# ---------- Image preprocessing ----------
def crop_image(src: Path, dst: Path, box):
    img = Image.open(src).convert("RGB")
    img.crop(box).save(dst, quality=95)

if (ASSET_DIR / "dashboard_current_full.png").exists():
    crop_image(ASSET_DIR / "dashboard_current_full.png", ASSET_DIR / "dashboard_top_crop.png", (0, 0, 1600, 900))
    crop_image(ASSET_DIR / "dashboard_current_full.png", ASSET_DIR / "dashboard_table_crop.png", (0, 760, 1600, 1670))
if (ASSET_DIR / "detail_current.png").exists():
    crop_image(ASSET_DIR / "detail_current.png", ASSET_DIR / "detail_crop.png", (0, 0, 1600, 900))
if (ASSET_DIR / "wiki_current.png").exists():
    crop_image(ASSET_DIR / "wiki_current.png", ASSET_DIR / "wiki_crop.png", (0, 0, 1000, 900))

# Mini chart image for data summary slide
W, H = 1100, 620
chart = Image.new("RGB", (W, H), "#08111f")
d = ImageDraw.Draw(chart)
try:
    font_b = ImageFont.truetype("malgunbd.ttf", 34)
    font = ImageFont.truetype("malgun.ttf", 24)
    font_s = ImageFont.truetype("malgun.ttf", 19)
except Exception:
    font_b = font = font_s = None

d.text((34, 24), "현재 데이터 스냅샷", fill="#F8FAFC", font=font_b)
palette = ["#2DD4BF", "#8BA8FF", "#A78BFA", "#FBBF24", "#FB7185"]

def draw_bars(counter, title, x, y, max_items=5):
    d.text((x, y), title, fill="#BAE6FD", font=font)
    items = counter.most_common(max_items)
    maxv = max([v for _, v in items] or [1])
    yy = y + 46
    for i, (label, value) in enumerate(items):
        label = str(label)
        d.text((x, yy), label[:25], fill="#D8E7FF", font=font_s)
        bar_x = x + 270
        bar_w = int(210 * value / maxv)
        d.rounded_rectangle([bar_x, yy + 5, bar_x + 210, yy + 20], radius=7, fill="#24324A")
        d.rounded_rectangle([bar_x, yy + 5, bar_x + bar_w, yy + 20], radius=7, fill=palette[i % len(palette)])
        d.text((bar_x + 230, yy - 2), str(value), fill="#FFFFFF", font=font_s)
        yy += 38

draw_bars(themes, "Theme 분포", 34, 95)
draw_bars(countries, "국가별 후보군", 560, 95)
draw_bars(clusters, "상위 Cluster", 34, 360)
draw_bars(statuses, "Filter 결과", 560, 360)
chart.save(ASSET_DIR / "data_snapshot.png", quality=95)

# ---------- PPT helpers ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

NAVY = RGBColor(5, 11, 24)
PANEL = RGBColor(17, 27, 43)
PANEL2 = RGBColor(24, 38, 60)
TEXT = RGBColor(244, 247, 251)
MUTED = RGBColor(176, 196, 222)
ACCENT = RGBColor(45, 212, 191)
BLUE = RGBColor(139, 168, 255)
GOLD = RGBColor(251, 191, 36)
RED = RGBColor(251, 113, 133)


def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, x, y, w, h, size=20, bold=False, color=TEXT, align=None, font="Malgun Gothic", margin=0.08):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_title(slide, title, subtitle=None, dark=True):
    textbox(slide, "SKBP Pipeline Finder", 0.55, 0.28, 4, 0.3, 11, True, ACCENT)
    textbox(slide, title, 0.55, 0.65, 9.8, 0.65, 34, True, TEXT)
    if subtitle:
        textbox(slide, subtitle, 0.58, 1.25, 10.6, 0.4, 15, False, MUTED)


def rounded_rect(slide, x, y, w, h, fill=PANEL, line=RGBColor(45, 60, 85), radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def metric(slide, label, value, x, y, w=2.25):
    rounded_rect(slide, x, y, w, 1.05, PANEL2)
    textbox(slide, label, x+0.16, y+0.15, w-0.3, 0.25, 10, True, MUTED)
    textbox(slide, value, x+0.16, y+0.43, w-0.3, 0.45, 23, True, TEXT)


def bullet_list(slide, items, x, y, w, h, size=15, color=TEXT, gap=0.1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return tb


def add_picture_fit(slide, path, x, y, w, h):
    path = str(path)
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # fit width
        pic_w = w; pic_h = w / img_ratio
        px = x; py = y + (h - pic_h) / 2
    else:
        pic_h = h; pic_w = h * img_ratio
        px = x + (w - pic_w) / 2; py = y
    return slide.shapes.add_picture(path, Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))


def footer(slide, idx):
    textbox(slide, f"{idx} / 11", 12.25, 7.12, 0.55, 0.2, 8, False, RGBColor(116, 139, 168), align=PP_ALIGN.RIGHT)

# ---------- Slides ----------
# 1
s = prs.slides.add_slide(blank); set_bg(s)
hero_visual = ASSET_DIR / "gen_hero_cockpit.png"
if hero_visual.exists():
    add_picture_fit(s, hero_visual, 0, 0, 13.333, 7.5)
textbox(s, "SKBP Pipeline Finder", 0.55, 0.55, 5.6, 0.45, 18, True, ACCENT)
textbox(s, "PreC Pipeline\nShortlisting Dashboard", 0.55, 1.12, 7.5, 1.6, 43, True, TEXT)
textbox(s, "GPT 조사 결과를 JSON 단일 원본으로 모아 후보 비교·근거 검토·Wiki 재사용까지 지원", 0.6, 2.9, 8.55, 0.62, 17, False, MUTED)
metric(s, "분석 레코드", f"{len(records)}건", 0.62, 4.1, 1.95)
metric(s, "PASS/SELECT", f"{pass_select}건", 2.78, 4.1, 1.95)
metric(s, "Target Relevance", f"{avg_target:.1f}/3", 4.94, 4.1, 1.95)
metric(s, "Wiki Graph", "477 nodes", 7.10, 4.1, 2.05)
if not hero_visual.exists():
    rounded_rect(s, 9.95, 0.55, 2.55, 5.1, RGBColor(8, 17, 31))
    add_picture_fit(s, ASSET_DIR / "dashboard_top_crop.png", 10.08, 0.86, 2.28, 4.2)
textbox(s, "전사 설명회 초안 · 7분", 0.62, 6.55, 4.5, 0.35, 14, True, GOLD)
footer(s, 1)

# 2
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "왜 만들었나: ‘후보 검토’의 병목을 줄이는 도구")
for i, (head, body, col) in enumerate([
    ("Before", "리포트·출처·점수가 흩어져\n후보 비교가 느림", RED),
    ("Need", "같은 기준으로 빠르게 선별하고\n점수 근거를 추적", GOLD),
    ("After", "JSON 대시보드에서\n비교·필터·근거 확인", ACCENT),
]):
    x = 0.75 + i*4.05
    rounded_rect(s, x, 1.85, 3.55, 2.6, PANEL2)
    textbox(s, head, x+0.25, 2.05, 2.2, 0.45, 24, True, col)
    textbox(s, body, x+0.25, 2.75, 2.95, 1.0, 17, False, TEXT)
textbox(s, "핵심 메시지", 0.78, 5.1, 2.5, 0.35, 18, True, ACCENT)
bullet_list(s, ["전략 후보 탐색을 ‘문서 읽기’에서 ‘데이터 기반 의사결정’으로 전환", "Fast Triage → Full Scout → Dashboard/Wiki로 이어지는 반복 가능한 프로세스", "출처·점수·불확실성을 함께 남겨 사후 검토와 인수인계가 쉬움"], 0.9, 5.55, 11.7, 1.1, 15)
footer(s, 2)

# 3
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "제작 과정: GPT 결과를 운영 가능한 파이프라인으로 구조화")
steps = [
    ("1", "GPT 1\nFast Triage", "여러 asset을 SELECT / REJECT / N/A로 1차 선별"),
    ("2", "GPT 2\nFull Scout", "선정 후보를 v3.1 rubric에 맞춰 깊게 조사"),
    ("3", "JSON Schema", "pipeline-records.json에 동일 구조로 저장"),
    ("4", "Dashboard", "필터·정렬·차트·상세 근거 검토"),
    ("5", "Wiki Export", "Obsidian/Markdown note와 graph로 확장"),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.55 + i*2.55
    rounded_rect(s, x, 2.0, 2.15, 2.25, PANEL2)
    textbox(s, num, x+0.18, 2.15, 0.45, 0.45, 22, True, ACCENT)
    textbox(s, title, x+0.22, 2.72, 1.75, 0.65, 18, True, TEXT)
    textbox(s, desc, x+0.22, 3.45, 1.72, 0.55, 11, False, MUTED)
    if i < 4:
        textbox(s, "→", x+2.22, 2.85, 0.35, 0.35, 26, True, BLUE)
workflow_visual = ASSET_DIR / "gen_workflow_strip.png"
if workflow_visual.exists():
    add_picture_fit(s, workflow_visual, 0.75, 4.65, 11.8, 1.25)
rounded_rect(s, 0.75, 6.0, 11.8, 0.42, RGBColor(8, 17, 31))
textbox(s, "운영 원칙: JSON이 단일 원본(Single Source of Truth)이고, 대시보드와 Wiki는 JSON에서 생성되는 산출물입니다.", 0.95, 6.09, 11.2, 0.18, 12, True, TEXT)
footer(s, 3)

# 4
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "어떻게 돌아가나: 단순하지만 확장 가능한 아키텍처")
# architecture boxes
arch = [
    ("브라우저 UI", "index/detail/wiki_view\nES module + CSS", 0.7, 2.0, BLUE),
    ("FastAPI", "main.py\nREST + static files", 3.45, 2.0, ACCENT),
    ("JSON Store", "json/pipeline-records.json\nschema + rubric", 6.2, 2.0, GOLD),
    ("Exports", "obsidian/\nskbp_pipeline_wiki/", 8.95, 2.0, RGBColor(167,139,250)),
    ("LLM Agent", "OpenRouter API\nJSON + Wiki retrieval", 6.2, 4.65, RED),
]
for title, body, x, y, col in arch:
    rounded_rect(s, x, y, 2.25, 1.25, PANEL2)
    textbox(s, title, x+0.18, y+0.2, 1.8, 0.3, 16, True, col)
    textbox(s, body, x+0.18, y+0.58, 1.85, 0.45, 11, False, TEXT)
for x1,y1,x2,y2 in [(2.95,2.62,3.35,2.62),(5.7,2.62,6.1,2.62),(8.45,2.62,8.85,2.62),(7.3,3.3,7.3,4.55)]:
    line=s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)); line.line.color.rgb=ACCENT; line.line.width=Pt(2)
textbox(s, "핵심 엔드포인트", 0.82, 5.35, 2.5, 0.3, 16, True, ACCENT)
bullet_list(s, ["GET /api/records, PUT/POST /api/records: 후보 데이터 조회·저장", "POST /api/wiki/export, /api/markdown/export: Markdown/Wiki 재생성", "POST /api/chat/stream: 선택 asset 또는 대시보드 맥락으로 AI 질의"], 1.0, 5.78, 11, 0.8, 13)
footer(s, 4)

# 5
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "대시보드 한눈에 보기: 후보군을 숫자와 차트로 즉시 파악")
add_picture_fit(s, ASSET_DIR / "dashboard_top_crop.png", 0.55, 1.55, 8.8, 5.2)
rounded_rect(s, 9.65, 1.62, 3.1, 5.05, PANEL2)
textbox(s, "화면 구성", 9.9, 1.9, 2.4, 0.35, 19, True, ACCENT)
bullet_list(s, ["총 분석 건수·PASS/SELECT·평균 TR 등 KPI", "Target/Theme/국가/Filter 분포 차트", "Priority Watch: 상위 후보 즉시 진입", "검색·Stage·Theme·Cluster·국가·적응증 필터", "테이블 정렬, 페이지 크기, 컬럼 설정, Excel export"], 9.95, 2.45, 2.45, 3.3, 13)
footer(s, 5)

# 6
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "현재 데이터: 32개 후보를 5개 관점으로 비교")
add_picture_fit(s, ASSET_DIR / "data_snapshot.png", 0.65, 1.45, 7.4, 4.2)
metric(s, "분석 건수", str(len(records)), 8.55, 1.55, 1.75)
metric(s, "PASS/SELECT", f"{pass_select}/{len(records)}", 10.55, 1.55, 1.75)
metric(s, "평균 총점", f"{avg_total:.1f}/21", 8.55, 2.85, 1.75)
metric(s, "평균 TR", f"{avg_target:.1f}/3", 10.55, 2.85, 1.75)
rounded_rect(s, 8.55, 4.25, 3.75, 1.75, PANEL2)
textbox(s, "해석 포인트", 8.78, 4.48, 2.5, 0.3, 16, True, ACCENT)
bullet_list(s, ["E/I Balance와 Neuroimmune 중심으로 분포", "중국·미국·한국 후보가 주요 비중", "Ion Channel, Cytokine/교세포 등 SKBP 관심 cluster로 재분류"], 8.8, 4.9, 3.15, 0.85, 11)
footer(s, 6)

# 7
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "판단 기준: 점수만이 아니라 ‘왜 그 점수인가’를 남김")
criteria = ["Target Relevance", "Competitive Landscape", "MoA Validity", "Platform Attractiveness", "Expansion Potential", "Data Maturity", "Marketability"]
for i, c in enumerate(criteria):
    x = 0.7 + (i%4)*3.05
    y = 1.65 + (i//4)*1.25
    rounded_rect(s, x, y, 2.55, 0.85, PANEL2)
    textbox(s, c, x+0.18, y+0.18, 2.15, 0.25, 13, True, TEXT)
    textbox(s, "0–3점 + Evidence Type", x+0.18, y+0.5, 2.1, 0.2, 10, False, MUTED)
rounded_rect(s, 0.8, 4.65, 5.6, 1.25, RGBColor(8, 17, 31))
textbox(s, "PASS / REVIEW / FAIL 게이트", 1.05, 4.9, 3.2, 0.3, 18, True, ACCENT)
textbox(s, "예: Total ≥14, TR ≥3, MOA ≥2, Data ≥2\n+ hard blocker 없음 → PASS 후보", 1.05, 5.25, 5.0, 0.55, 13, False, TEXT)
rounded_rect(s, 6.8, 4.65, 5.6, 1.25, RGBColor(8, 17, 31))
textbox(s, "Audit label", 7.05, 4.9, 3.2, 0.3, 18, True, GOLD)
textbox(s, "E0–E4 evidence type, why_not_higher,\ninvestigation_note, source URL 보존", 7.05, 5.25, 4.8, 0.55, 13, False, TEXT)
footer(s, 7)

# 8
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "상세 화면: 원문 리포트와 Score 근거를 같은 화면에서 검토")
add_picture_fit(s, ASSET_DIR / "detail_crop.png", 0.55, 1.5, 8.9, 5.25)
rounded_rect(s, 9.75, 1.7, 2.75, 4.7, PANEL2)
textbox(s, "Detail의 가치", 9.98, 1.98, 2.2, 0.3, 18, True, ACCENT)
bullet_list(s, ["좌측 Mini TOC로 긴 GPT 리포트 탐색", "중앙: Markdown 원문 리포트", "우측: criteria별 점수·근거·source 링크", "AI Agent로 현재 asset 맥락 질문"], 10.0, 2.45, 2.2, 2.35, 12)
agent_visual = ASSET_DIR / "gen_agent_evidence.png"
if agent_visual.exists():
    add_picture_fit(s, agent_visual, 9.98, 5.15, 2.28, 0.95)
footer(s, 8)

# 9
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "Wiki/Obsidian 레이어: Dashboard 밖 지식베이스")
add_picture_fit(s, ASSET_DIR / "wiki_crop.png", 0.7, 1.45, 4.1, 5.35)
rounded_rect(s, 5.15, 1.6, 3.35, 1.05, PANEL2)
textbox(s, "자동 생성 노트", 5.4, 1.85, 2.2, 0.3, 18, True, ACCENT)
textbox(s, "Assets / Companies / Targets / Indications / Scorecards", 5.4, 2.25, 2.55, 0.22, 11, False, TEXT)
rounded_rect(s, 5.15, 3.0, 3.35, 1.05, PANEL2)
textbox(s, "대시보드 노트", 5.4, 3.25, 2.2, 0.3, 18, True, BLUE)
textbox(s, "Asset Index, By Target, By Theme, Evidence Gaps", 5.4, 3.65, 2.55, 0.22, 11, False, TEXT)
rounded_rect(s, 5.15, 4.4, 3.35, 1.05, PANEL2)
textbox(s, "Graph Export", 5.4, 4.65, 2.2, 0.3, 18, True, GOLD)
textbox(s, "nodes.csv / edges.csv / graph.json — 477 nodes, 782 edges", 5.4, 5.05, 2.6, 0.22, 11, False, TEXT)
rounded_rect(s, 8.95, 1.6, 3.25, 3.85, RGBColor(8, 17, 31))
textbox(s, "운영 방식", 9.2, 1.9, 2.2, 0.3, 18, True, ACCENT)
bullet_list(s, ["JSON 수정 → export script 실행", "Markdown vault는 산출물이므로 재생성 가능", "LLM Agent가 dashboard JSON + Wiki note를 함께 검색", "회의 후 Obsidian에서 연결 지식으로 계속 활용"], 9.2, 2.35, 2.5, 2.1, 12)
footer(s, 9)

# 10
s = prs.slides.add_slide(blank); set_bg(s); add_title(s, "전사 데모 7분 구성안")
agenda = [
    ("0:00–0:40", "문제 정의", "후보 검토가 왜 오래 걸렸는지"),
    ("0:40–1:30", "제작 과정", "GPT1/2, JSON schema, rubric"),
    ("1:30–3:10", "Dashboard Demo", "KPI·차트·필터·Priority Watch"),
    ("3:10–4:40", "Detail Demo", "원문, 점수근거, 출처, AI Agent"),
    ("4:40–5:40", "Wiki Layer", "Obsidian export와 graph"),
    ("5:40–6:40", "운영/확장", "데이터 갱신, 배포, DB 전환"),
    ("6:40–7:00", "마무리", "의사결정 속도와 감사 가능성"),
]
for i, (time, title, desc) in enumerate(agenda):
    y = 1.55 + i*0.72
    textbox(s, time, 0.85, y, 1.3, 0.25, 13, True, ACCENT)
    textbox(s, title, 2.25, y, 2.0, 0.25, 15, True, TEXT)
    textbox(s, desc, 4.35, y, 7.5, 0.25, 13, False, MUTED)
footer(s, 10)

# 11
s = prs.slides.add_slide(blank); set_bg(s)
add_title(s, "마무리: 이 대시보드의 핵심 특징")
for i, (title, body, col) in enumerate([
    ("비교 가능성", "후보를 같은 점수체계와 필터로 비교", ACCENT),
    ("감사 가능성", "점수·근거·출처·불확실성까지 보존", GOLD),
    ("확장 가능성", "JSON → Dashboard → Wiki → Agent로 재사용", BLUE),
]):
    x = 0.75 + i*4.05
    rounded_rect(s, x, 1.9, 3.55, 2.2, PANEL2)
    textbox(s, title, x+0.25, 2.2, 2.6, 0.4, 23, True, col)
    textbox(s, body, x+0.25, 2.85, 2.85, 0.6, 15, False, TEXT)
rounded_rect(s, 1.2, 5.1, 10.9, 0.85, RGBColor(8, 17, 31))
textbox(s, "한 줄 요약: 흩어진 GPT 조사 결과를 전사적으로 공유 가능한 ‘pipeline intelligence cockpit’으로 바꾼 프로젝트입니다.", 1.48, 5.38, 10.2, 0.28, 17, True, TEXT)
textbox(s, "Q&A", 5.45, 6.45, 2.4, 0.45, 24, True, ACCENT, align=PP_ALIGN.CENTER)
footer(s, 11)

prs.save(OUT)

# ---------- Speaker draft ----------
script = f"""# SKBP Pipeline Dashboard 7분 발표 원고 초안

## 1. 오프닝 (0:00–0:40)
안녕하세요. 오늘은 SKBP Pipeline Finder, 즉 PreC pipeline 후보를 빠르게 선별하고 근거까지 추적할 수 있도록 만든 내부 대시보드를 소개드리겠습니다. 이 프로젝트의 목적은 GPT로 조사한 후보 리포트를 단순 문서로 끝내지 않고, 같은 기준으로 비교 가능한 데이터 자산으로 전환하는 것입니다. 현재 JSON 기준으로 총 {len(records)}개 후보가 들어 있고, PASS 또는 SELECT로 볼 수 있는 후보는 {pass_select}개입니다.

## 2. 제작 배경과 문제 (0:40–1:30)
기존에는 회사별·asset별 GPT 리포트가 각각 존재해 후보 간 비교가 어렵고, 점수의 근거와 출처를 다시 찾는 시간이 많이 들었습니다. 그래서 세 가지를 해결하려고 했습니다. 첫째, 후보를 같은 rubric으로 비교할 것. 둘째, 왜 그 점수가 나왔는지 evidence trail을 남길 것. 셋째, 대시보드와 Wiki로 공유 가능한 형태를 만들 것입니다.

## 3. 제작 과정 (1:30–2:20)
흐름은 GPT 1 Fast Triage, GPT 2 Full Scout, JSON Schema 저장, Dashboard 표시, Wiki Export 순서입니다. GPT 1은 여러 asset을 SELECT, REJECT, N/A로 빠르게 선별합니다. SELECT 후보는 GPT 2에서 Target Relevance, MoA Validity, Data Maturity, Marketability 등 7개 기준으로 심층 조사합니다. 결과는 `json/pipeline-records.json`에 저장되고, 이 JSON이 단일 원본입니다.

## 4. 대시보드 소개 (2:20–3:40)
메인 화면에서는 총 분석 건수, PASS/SELECT 비율, 평균 Target Relevance, 국가 수 같은 KPI를 먼저 보여줍니다. 아래에는 Target Relevance 분포, Theme 분포, 국가별 후보군, Pipeline Filter, Score Profile, Priority Watch가 있습니다. 필터는 Stage, Theme, Cluster, 국가, indication, Pipeline Filter 기준으로 걸 수 있고, 테이블은 정렬·컬럼 설정·Excel export가 가능합니다. 즉 회의 중에도 “Neuroimmune 중에서 한국 회사 후보만 보자” 같은 질문에 바로 대응할 수 있습니다.

## 5. 평가 기준과 Hard Filter (3:40–4:40)
점수 체계는 7개 criteria 각각 0~3점입니다. 단순 총점만 보는 것이 아니라 Target Relevance, MoA, Data Maturity 같은 핵심 게이트와 hard blocker를 함께 봅니다. 예를 들어 Total 14점 이상, TR 3점, MoA 2점 이상, Data 2점 이상이고 명확한 blocker가 없으면 PASS 후보로 볼 수 있습니다. 또 각 점수에는 evidence type, why_not_higher, investigation note, source URL이 함께 남아 사후 검증이 가능합니다.

## 6. 상세 화면과 AI Agent (4:40–5:40)
후보 하나를 클릭하면 상세 화면으로 들어갑니다. 좌측에는 긴 리포트의 Mini TOC, 중앙에는 GPT 원문 리포트, 우측에는 score별 판단 근거와 출처가 나옵니다. 또한 Asset Evidence Agent가 있어서 현재 asset의 JSON과 Wiki note를 맥락으로 Target fit, Marketability, Evidence gap, Competitor risk를 질문할 수 있습니다. 이 기능은 단순 채팅이 아니라 현재 후보 데이터와 Wiki 검색 결과를 함께 참고하도록 설계되어 있습니다.

## 7. Wiki/Obsidian 레이어 (5:40–6:25)
대시보드와 별도로 `skbp_pipeline_wiki`가 자동 생성됩니다. 여기에는 asset, company, target, indication, competitor, scorecard, theme, cluster별 note가 있고, dashboard note와 graph export도 포함됩니다. 현재 Wiki README 기준 graph nodes는 477개, edges는 782개입니다. 따라서 대시보드에서 본 내용을 Obsidian 같은 지식베이스에서도 연결 관계로 볼 수 있습니다.

## 8. 운영 방식과 확장 (6:25–6:50)
운영은 간단합니다. JSON을 저장하고, 필요하면 Markdown/Wiki export를 다시 실행합니다. 현재는 로컬 JSON 파일 기반이지만 FastAPI 구조라 Render, Railway 같은 Python web service로 올릴 수 있고, 여러 명이 동시에 쓰려면 다음 단계에서 SQLite나 Postgres로 전환하면 됩니다.

## 9. 클로징 (6:50–7:00)
정리하면 이 대시보드는 GPT 조사 결과를 전사적으로 공유 가능한 pipeline intelligence cockpit으로 바꾼 프로젝트입니다. 후보 비교 속도를 높이고, 점수의 근거를 남기며, Wiki와 AI Agent로 재사용성을 확보한 것이 핵심 특징입니다. 감사합니다.
"""
SCRIPT_OUT.write_text(script, encoding="utf-8")
print(f"PPTX: {OUT}")
print(f"SCRIPT: {SCRIPT_OUT}")
print(f"Slides: {len(prs.slides)}")
