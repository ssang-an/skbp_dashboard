from __future__ import annotations

import base64
import hashlib
import json
import os
import re
import shutil
import subprocess
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Literal

import fitz
import requests
import urllib3
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, ValidationError

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

DEFAULT_OPENROUTER_BASE_URL = "https://openrouter.ai/api/v1"
DEFAULT_PDF_PARSER_ENGINE = "cloudflare-ai"
PAID_PDF_PARSER_ENGINE = "mistral-ocr"
FILTER3_DOCUMENT_TEXT_LIMIT = 80_000
FILTER3_ANALYSIS_MAX_TOKENS = 3_000
DOCUMENT_KEYWORDS = re.compile(
    r"\b(in[\s-]?vitro|in[\s-]?vivo|animal\s+model|mouse|mice|rat|pk|pd|pk/pd|"
    r"efficacy|potency|admet|tox(?:icity)?|pharmacokinetic|pharmacodynamic)\b|"
    r"(유효성|효능|동물\s*모델|독성|약동학|약력학)",
    re.IGNORECASE,
)
REPEATED_CHARACTER_RUN = re.compile(r"(.)\1{7,}", re.DOTALL)
JSON_FENCE = re.compile(r"^```(?:json)?\s*|\s*```$", re.IGNORECASE)


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def env_bool(name: str, default: bool) -> bool:
    value = os.getenv(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def openrouter_base_url() -> str:
    configured = (
        os.getenv("OPENROUTER_BASE_URL")
        or os.getenv("OPENROUTER_API_URL")
        or DEFAULT_OPENROUTER_BASE_URL
    ).strip().rstrip("/")
    if configured.endswith("/chat/completions"):
        return configured[: -len("/chat/completions")]
    return configured


def openrouter_chat_url() -> str:
    return f"{openrouter_base_url()}/chat/completions"


def openrouter_headers(api_key: str) -> dict[str, str]:
    return {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": os.getenv("OPENROUTER_SITE_URL", "http://127.0.0.1:8011"),
        "X-Title": os.getenv("OPENROUTER_APP_TITLE", "SKBP Pipeline Finder"),
    }


def sha256_file(file_path: Path) -> str:
    digest = hashlib.sha256()
    with file_path.open("rb") as source:
        for block in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def labeled_text(pages_or_slides: list[dict[str, Any]], label: str) -> str:
    sections = []
    for item in pages_or_slides:
        text = str(item.get("text") or "").strip()
        if not text:
            continue
        number = item.get("number")
        heading = f"[{label} {number}]" if number is not None else f"[{label}]"
        sections.append(f"{heading}\n{text}")
    return "\n\n".join(sections)


def extract_pdf_native(file_path: Path) -> dict[str, Any]:
    pages: list[dict[str, Any]] = []
    with fitz.open(file_path) as document:
        total_pages = document.page_count
        for index, page in enumerate(document, 1):
            pages.append({"number": index, "text": page.get_text("text").strip()})
    return {
        "method": "native_pdf_text",
        "pages_or_slides": pages,
        "page_or_slide_count": total_pages,
        "parsed_text": labeled_text(pages, "PAGE"),
    }


def _shape_text(shape: Any) -> list[str]:
    values: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = str(getattr(shape, "text", "") or "").strip()
        if text:
            values.append(text)
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                values.append(" | ".join(cells))
    if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            values.extend(_shape_text(child))
    return values


def extract_pptx_native(file_path: Path) -> dict[str, Any]:
    presentation = Presentation(str(file_path))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_shape_text(shape))
        slides.append({"number": index, "text": "\n".join(texts).strip()})
    return {
        "method": "native_pptx_text",
        "pages_or_slides": slides,
        "page_or_slide_count": len(slides),
        "parsed_text": labeled_text(slides, "SLIDE"),
    }


def extraction_quality(
    pages_or_slides: list[dict[str, Any]],
    *,
    minimum_characters: int,
) -> dict[str, Any]:
    texts = [str(item.get("text") or "") for item in pages_or_slides]
    combined = "\n".join(texts)
    non_whitespace = re.sub(r"\s+", "", combined)
    character_count = len(non_whitespace)
    corrupt_count = sum(
        1
        for character in non_whitespace
        if character == "\ufffd" or (ord(character) < 32 and character not in "\n\r\t")
    )
    repeated_count = sum(len(match.group(0)) for match in REPEATED_CHARACTER_RUN.finditer(combined))
    text_units = sum(1 for text in texts if re.search(r"\S", text))
    total_units = len(texts)
    text_unit_ratio = text_units / total_units if total_units else 0.0
    corrupted_character_ratio = corrupt_count / max(character_count, 1)
    repeated_character_ratio = repeated_count / max(character_count, 1)
    keyword_hits = sorted({match.group(0).lower() for match in DOCUMENT_KEYWORDS.finditer(combined)})

    score = 0
    if character_count >= minimum_characters:
        score += 2
    if corrupted_character_ratio <= 0.12:
        score += 1
    if repeated_character_ratio <= 0.25:
        score += 1
    if text_unit_ratio >= 0.25:
        score += 1
    if keyword_hits:
        score += 1

    sufficient = (
        character_count >= minimum_characters
        and corrupted_character_ratio <= 0.20
        and repeated_character_ratio <= 0.35
        and text_unit_ratio >= 0.25
        and bool(keyword_hits)
        and score >= 5
    )
    reasons = []
    if character_count < minimum_characters:
        reasons.append("insufficient_characters")
    if corrupted_character_ratio > 0.20:
        reasons.append("high_corrupted_character_ratio")
    if repeated_character_ratio > 0.35:
        reasons.append("abnormal_character_repetition")
    if text_unit_ratio < 0.25:
        reasons.append("low_text_page_ratio")
    if not keyword_hits:
        reasons.append("no_filter3_keywords")

    return {
        "sufficient": sufficient,
        "quality_score": score,
        "character_count": character_count,
        "corrupted_character_ratio": round(corrupted_character_ratio, 5),
        "repeated_character_ratio": round(repeated_character_ratio, 5),
        "text_unit_ratio": round(text_unit_ratio, 5),
        "text_units": text_units,
        "total_units": total_units,
        "keyword_hits": keyword_hits[:30],
        "reasons": reasons,
    }


def find_libreoffice() -> Path | None:
    configured = os.getenv("LIBREOFFICE_PATH")
    candidates = [
        Path(configured) if configured else None,
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ]
    command = shutil.which("soffice") or shutil.which("libreoffice")
    if command:
        candidates.insert(0, Path(command))
    return next((candidate for candidate in candidates if candidate and candidate.is_file()), None)


def convert_office_to_pdf(file_path: Path) -> dict[str, Any]:
    soffice = find_libreoffice()
    if soffice is None:
        return {"status": "unavailable", "pdf_path": None, "error": "LibreOffice was not found."}
    output_dir = file_path.parent / "converted"
    output_dir.mkdir(parents=True, exist_ok=True)
    completed = subprocess.run(
        [
            str(soffice),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(file_path),
        ],
        cwd=str(file_path.parent),
        capture_output=True,
        text=True,
        timeout=180,
        check=False,
    )
    output_path = output_dir / f"{file_path.stem}.pdf"
    if completed.returncode == 0 and output_path.is_file():
        return {"status": "converted", "pdf_path": str(output_path), "error": None}
    detail = (completed.stderr or completed.stdout or "LibreOffice conversion failed.").strip()
    return {"status": "failed", "pdf_path": None, "error": detail[:500]}


def _file_annotations(response_data: dict[str, Any]) -> list[dict[str, Any]]:
    choices = response_data.get("choices") if isinstance(response_data, dict) else None
    message = choices[0].get("message") if isinstance(choices, list) and choices else None
    annotations = message.get("annotations") if isinstance(message, dict) else None
    if isinstance(annotations, list):
        return [item for item in annotations if isinstance(item, dict) and item.get("type") == "file"]
    error = response_data.get("error") if isinstance(response_data, dict) else None
    metadata = error.get("metadata") if isinstance(error, dict) else None
    annotations = metadata.get("file_annotations") if isinstance(metadata, dict) else None
    return [item for item in annotations or [] if isinstance(item, dict) and item.get("type") == "file"]


def _sanitized_annotations(annotations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    sanitized = []
    for annotation in annotations:
        file_data = annotation.get("file") if isinstance(annotation.get("file"), dict) else {}
        content = [
            {"type": "text", "text": str(part.get("text") or "")}
            for part in file_data.get("content") or []
            if isinstance(part, dict) and part.get("type") == "text" and part.get("text")
        ]
        sanitized.append(
            {
                "type": "file",
                "file": {
                    "hash": str(file_data.get("hash") or ""),
                    "name": str(file_data.get("name") or ""),
                    "content": content,
                },
            }
        )
    return sanitized


def _annotation_text(annotations: list[dict[str, Any]]) -> tuple[str, str]:
    texts: list[str] = []
    file_hash = ""
    for annotation in annotations:
        file_data = annotation.get("file") if isinstance(annotation.get("file"), dict) else {}
        file_hash = file_hash or str(file_data.get("hash") or "")
        for part in file_data.get("content") or []:
            if isinstance(part, dict) and part.get("type") == "text" and part.get("text"):
                texts.append(str(part["text"]))
    return "\n\n".join(texts).strip(), file_hash


def parse_pdf_with_openrouter(
    file_path: Path,
    filename: str,
    *,
    engine: str,
    api_key: str,
    model: str,
) -> dict[str, Any]:
    data_url = "data:application/pdf;base64," + base64.b64encode(file_path.read_bytes()).decode("ascii")
    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Parse this PDF faithfully. Return only PARSING_COMPLETE. "
                            "Do not summarize or infer; the parsed file annotation is the required output."
                        ),
                    },
                    {
                        "type": "file",
                        "file": {"filename": filename, "file_data": data_url},
                    },
                ],
            }
        ],
        "plugins": [{"id": "file-parser", "pdf": {"engine": engine}}],
        "temperature": 0,
        "max_tokens": 30,
    }
    try:
        response = requests.post(
            openrouter_chat_url(),
            headers=openrouter_headers(api_key),
            json=payload,
            timeout=int(os.getenv("OPENROUTER_PDF_TIMEOUT_SECONDS", "180")),
            verify=False,
        )
        data = response.json()
    except Exception as exc:
        return {
            "status": "failed",
            "parser_engine": engine,
            "parsed_text": "",
            "file_hash": "",
            "file_annotations": [],
            "parsed_at": utc_now(),
            "error": str(exc)[:700],
        }

    annotations = _file_annotations(data)
    sanitized = _sanitized_annotations(annotations)
    parsed_text, file_hash = _annotation_text(sanitized)
    if parsed_text:
        return {
            "status": "completed",
            "parser_engine": engine,
            "parsed_text": parsed_text,
            "file_hash": file_hash,
            "file_annotations": sanitized,
            "parsed_at": utc_now(),
            "error": None,
        }
    error = data.get("error") if isinstance(data, dict) else None
    return {
        "status": "failed",
        "parser_engine": engine,
        "parsed_text": "",
        "file_hash": file_hash,
        "file_annotations": sanitized,
        "parsed_at": utc_now(),
        "error": json.dumps(error or data, ensure_ascii=False)[:700],
    }


class EvidenceExcerpt(BaseModel):
    model_config = ConfigDict(extra="forbid")

    page_or_slide: int | None
    location_status: Literal["located", "page_not_available_from_parser"]
    text: str = Field(min_length=1, max_length=800)


class CriterionJudgment(BaseModel):
    model_config = ConfigDict(extra="forbid")

    verdict: Literal["true", "false", "unknown"]
    rationale: str = Field(min_length=1, max_length=800)
    evidence: list[EvidenceExcerpt] = Field(default_factory=list, max_length=8)


class Filter3DocumentJudgment(BaseModel):
    model_config = ConfigDict(extra="forbid")

    target_indication: CriterionJudgment
    small_molecule: CriterionJudgment
    ind_enabling: CriterionJudgment
    in_vivo_efficacy: CriterionJudgment
    in_vitro_efficacy: CriterionJudgment
    admet_at_least_25: CriterionJudgment
    admet_completed_count: int | None = Field(default=None, ge=0, le=50)
    matched_indication: str | None = Field(default=None, max_length=200)
    concise_summary: str = Field(min_length=1, max_length=1000)


_MODEL_PARAMETER_CACHE: dict[str, set[str]] = {}


def model_supported_parameters(model: str) -> set[str]:
    if model in _MODEL_PARAMETER_CACHE:
        return _MODEL_PARAMETER_CACHE[model]
    try:
        response = requests.get(
            f"{openrouter_base_url()}/models",
            timeout=30,
            verify=False,
        )
        response.raise_for_status()
        item = next(
            (candidate for candidate in response.json().get("data", []) if candidate.get("id") == model),
            None,
        )
        parameters = set(item.get("supported_parameters") or []) if isinstance(item, dict) else set()
    except Exception:
        parameters = set()
    _MODEL_PARAMETER_CACHE[model] = parameters
    return parameters


def _message_content(data: dict[str, Any]) -> str:
    try:
        content = data["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError):
        return ""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        return "\n".join(
            str(part.get("text") or "")
            for part in content
            if isinstance(part, dict) and part.get("type") == "text"
        )
    return ""


def _parse_judgment(content: str) -> Filter3DocumentJudgment:
    candidate = JSON_FENCE.sub("", content.strip())
    start = candidate.find("{")
    end = candidate.rfind("}")
    if start >= 0 and end > start:
        candidate = candidate[start : end + 1]
    return Filter3DocumentJudgment.model_validate_json(candidate)


def _normalized_evidence_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip().casefold()


def normalize_evidence_locations(
    judgment: Filter3DocumentJudgment,
    *,
    extraction_method: str,
    pages_or_slides: list[dict[str, Any]],
) -> Filter3DocumentJudgment:
    """Verify locations against source text instead of trusting model-supplied page numbers."""
    parser_without_pages = extraction_method.startswith("openrouter_file_parser")
    source_units = [
        (item.get("number"), _normalized_evidence_text(str(item.get("text") or "")))
        for item in pages_or_slides
        if _normalized_evidence_text(str(item.get("text") or ""))
    ]
    criterion_names = (
        "target_indication",
        "small_molecule",
        "ind_enabling",
        "in_vivo_efficacy",
        "in_vitro_efficacy",
        "admet_at_least_25",
    )
    for criterion_name in criterion_names:
        criterion = getattr(judgment, criterion_name)
        verified: list[EvidenceExcerpt] = []
        for evidence in criterion.evidence:
            if parser_without_pages:
                evidence.page_or_slide = None
                evidence.location_status = "page_not_available_from_parser"
                verified.append(evidence)
                continue

            excerpt = _normalized_evidence_text(evidence.text)
            matching_units = [
                number
                for number, source_text in source_units
                if excerpt and excerpt in source_text and isinstance(number, int)
            ]
            if not matching_units:
                # A paraphrase is not acceptable as a page-located quotation.
                continue
            if evidence.page_or_slide in matching_units:
                located_number = evidence.page_or_slide
            else:
                located_number = matching_units[0]
            evidence.page_or_slide = located_number
            evidence.location_status = "located"
            verified.append(evidence)
        criterion.evidence = verified
    return judgment


def analyze_filter3_with_deepseek(
    *,
    document_id: str,
    filename: str,
    extraction_method: str,
    pages_or_slides: list[dict[str, Any]],
    filter3_criteria: str,
    api_key: str,
    model: str,
) -> dict[str, Any]:
    schema = Filter3DocumentJudgment.model_json_schema()
    source_payload = {
        "document_id": document_id,
        "filename": filename,
        "extraction_method": extraction_method,
        "pages_or_slides": pages_or_slides,
        "filter3_criteria": filter3_criteria,
    }
    source_json = json.dumps(source_payload, ensure_ascii=False)
    if len(source_json) > FILTER3_DOCUMENT_TEXT_LIMIT:
        source_payload["pages_or_slides"] = pages_or_slides[:]
        remaining = FILTER3_DOCUMENT_TEXT_LIMIT
        trimmed = []
        for item in source_payload["pages_or_slides"]:
            text = str(item.get("text") or "")
            if remaining <= 0:
                break
            trimmed_text = text[:remaining]
            trimmed.append({**item, "text": trimmed_text})
            remaining -= len(trimmed_text)
        source_payload["pages_or_slides"] = trimmed
        source_json = json.dumps(source_payload, ensure_ascii=False)

    system_prompt = (
        "You are a document evidence extractor for SKBP Filter 3. "
        "Return factual judgments only from the supplied document text. "
        "Every verdict must be exactly true, false, or unknown. "
        "true means explicit supporting evidence. false means the document explicitly negates the criterion. "
        "Missing information is always unknown, never false. "
        "For in-vivo or in-vitro efficacy, merely saying an experiment was performed is not true; "
        "an explicit positive efficacy/activity result is required. "
        "Every evidence text must be a short verbatim excerpt copied from the supplied text, not a paraphrase. "
        "Preserve the supplied page or slide number in evidence. "
        "When extraction_method starts with openrouter_file_parser and no reliable page exists, "
        "use page_or_slide null and location_status page_not_available_from_parser. "
        "Do not invent page numbers, counts, results, or conclusions."
    )
    user_prompt = (
        "Analyze this document against the supplied Filter 3 criteria and output only one JSON object "
        "that conforms to the requested schema.\n\n"
        f"{source_json}"
    )
    parameters = model_supported_parameters(model)
    structured_supported = bool({"structured_outputs", "response_format"} & parameters)
    attempts: list[tuple[bool, int]] = (
        [(True, 2), (False, 2)] if structured_supported else [(False, 2)]
    )
    errors: list[str] = []

    for structured, retry_count in attempts:
        for _ in range(retry_count):
            payload: dict[str, Any] = {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                "temperature": 0,
                "max_tokens": FILTER3_ANALYSIS_MAX_TOKENS,
            }
            if structured:
                payload["response_format"] = {
                    "type": "json_schema",
                    "json_schema": {
                        "name": "filter3_document_judgment",
                        "strict": True,
                        "schema": schema,
                    },
                }
                payload["provider"] = {"require_parameters": True}
            else:
                payload["messages"][1]["content"] += (
                    "\n\nJSON Schema:\n" + json.dumps(schema, ensure_ascii=False)
                )
            try:
                response = requests.post(
                    openrouter_chat_url(),
                    headers=openrouter_headers(api_key),
                    json=payload,
                    timeout=int(os.getenv("OPENROUTER_ANALYSIS_TIMEOUT_SECONDS", "180")),
                    verify=False,
                )
                response.raise_for_status()
                data = response.json()
                judgment = _parse_judgment(_message_content(data))
                judgment = normalize_evidence_locations(
                    judgment,
                    extraction_method=extraction_method,
                    pages_or_slides=pages_or_slides,
                )
                return {
                    "status": "completed",
                    "model": model,
                    "structured_output_supported": structured_supported,
                    "structured_output_used": structured,
                    "analyzed_at": utc_now(),
                    "result": judgment.model_dump(mode="json"),
                    "error": None,
                }
            except (requests.RequestException, ValueError, ValidationError, KeyError) as exc:
                errors.append(str(exc)[:700])

    return {
        "status": "failed",
        "model": model,
        "structured_output_supported": structured_supported,
        "structured_output_used": False,
        "analyzed_at": utc_now(),
        "result": None,
        "error": " / ".join(errors[-4:])[:1500] or "DeepSeek analysis failed.",
    }


def parser_pages(parsed_text: str) -> list[dict[str, Any]]:
    return [
        {
            "number": None,
            "location_status": "page_not_available_from_parser",
            "text": parsed_text,
        }
    ]


def process_document(
    file_path: Path,
    filename: str,
    *,
    filter3_criteria: str,
    cached_parser: dict[str, Any] | None = None,
) -> dict[str, Any]:
    document_id = uuid.uuid4().hex
    file_sha256 = sha256_file(file_path)
    suffix = file_path.suffix.lower()
    minimum_characters = max(1, int(os.getenv("PDF_TEXT_MIN_CHARACTERS", "50")))
    api_key = os.getenv("OPENROUTER_API_KEY", "").strip()
    model = os.getenv("OPENROUTER_MODEL", "").strip()
    parser_enabled = env_bool("ENABLE_PDF_PARSER_FALLBACK", True)
    parser_engine = os.getenv("PDF_PARSER_ENGINE", DEFAULT_PDF_PARSER_ENGINE).strip() or DEFAULT_PDF_PARSER_ENGINE
    paid_ocr_enabled = env_bool("ENABLE_PAID_OCR", False)
    conversion = {"status": "not_required", "pdf_path": None, "error": None}
    parser_result: dict[str, Any] | None = None
    extraction: dict[str, Any]

    try:
        if suffix == ".pdf":
            extraction = extract_pdf_native(file_path)
        elif suffix == ".pptx":
            extraction = extract_pptx_native(file_path)
            conversion = convert_office_to_pdf(file_path)
        elif suffix == ".ppt":
            conversion = convert_office_to_pdf(file_path)
            if not conversion.get("pdf_path"):
                raise RuntimeError(conversion.get("error") or "PPT conversion failed.")
            extraction = extract_pdf_native(Path(str(conversion["pdf_path"])))
            extraction["method"] = "libreoffice_pdf_native_text"
        else:
            return {
                "document_id": document_id,
                "file_sha256": file_sha256,
                "status": "not_applicable",
                "filename": filename,
                "processed_at": utc_now(),
            }
    except Exception as exc:
        extraction = {
            "method": "native_extraction_failed",
            "pages_or_slides": [],
            "page_or_slide_count": 0,
            "parsed_text": "",
            "error": str(exc)[:700],
        }

    quality = extraction_quality(
        extraction.get("pages_or_slides") or [],
        minimum_characters=minimum_characters,
    )
    extraction["quality"] = quality
    pdf_for_parser = file_path if suffix == ".pdf" else (
        Path(str(conversion["pdf_path"])) if conversion.get("pdf_path") else None
    )

    if not quality["sufficient"] and parser_enabled and pdf_for_parser is not None:
        if (
            isinstance(cached_parser, dict)
            and cached_parser.get("status") == "completed"
            and cached_parser.get("parsed_text")
        ):
            parser_result = {**cached_parser, "cache_hit": True}
        elif not api_key:
            parser_result = {
                "status": "failed",
                "parser_engine": parser_engine,
                "parsed_text": "",
                "file_hash": "",
                "file_annotations": [],
                "parsed_at": utc_now(),
                "error": "OPENROUTER_API_KEY 설정 또는 유효성 확인 필요",
                "cache_hit": False,
            }
        elif not model:
            parser_result = {
                "status": "failed",
                "parser_engine": parser_engine,
                "parsed_text": "",
                "file_hash": "",
                "file_annotations": [],
                "parsed_at": utc_now(),
                "error": "OPENROUTER_MODEL is not set.",
                "cache_hit": False,
            }
        else:
            parser_result = parse_pdf_with_openrouter(
                pdf_for_parser,
                filename,
                engine=parser_engine,
                api_key=api_key,
                model=model,
            )
            parser_result["cache_hit"] = False

        if parser_result.get("parsed_text"):
            parsed_pages = parser_pages(str(parser_result["parsed_text"]))
            parsed_quality = extraction_quality(
                parsed_pages,
                minimum_characters=minimum_characters,
            )
            parser_result["quality"] = parsed_quality
            extraction = {
                "method": f"openrouter_file_parser_{parser_result.get('parser_engine')}",
                "pages_or_slides": parsed_pages,
                "page_or_slide_count": None,
                "parsed_text": str(parser_result["parsed_text"]),
                "quality": parsed_quality,
            }

        if (
            parser_result
            and not parser_result.get("quality", {}).get("sufficient")
            and parser_result.get("parser_engine") != PAID_PDF_PARSER_ENGINE
            and paid_ocr_enabled
            and api_key
            and model
        ):
            paid_result = parse_pdf_with_openrouter(
                pdf_for_parser,
                filename,
                engine=PAID_PDF_PARSER_ENGINE,
                api_key=api_key,
                model=model,
            )
            paid_result["cache_hit"] = False
            if paid_result.get("parsed_text"):
                paid_pages = parser_pages(str(paid_result["parsed_text"]))
                paid_result["quality"] = extraction_quality(
                    paid_pages,
                    minimum_characters=minimum_characters,
                )
                parser_result = paid_result
                extraction = {
                    "method": f"openrouter_file_parser_{PAID_PDF_PARSER_ENGINE}",
                    "pages_or_slides": paid_pages,
                    "page_or_slide_count": None,
                    "parsed_text": str(paid_result["parsed_text"]),
                    "quality": paid_result["quality"],
                }

    analysis: dict[str, Any]
    if not api_key:
        analysis = {
            "status": "failed",
            "model": model or None,
            "result": None,
            "analyzed_at": utc_now(),
            "error": "OPENROUTER_API_KEY 설정 또는 유효성 확인 필요",
        }
    elif not model:
        analysis = {
            "status": "failed",
            "model": None,
            "result": None,
            "analyzed_at": utc_now(),
            "error": "OPENROUTER_MODEL is not set.",
        }
    elif not extraction.get("parsed_text"):
        analysis = {
            "status": "failed",
            "model": model,
            "result": None,
            "analyzed_at": utc_now(),
            "error": "No extractable document text.",
        }
    else:
        analysis = analyze_filter3_with_deepseek(
            document_id=document_id,
            filename=filename,
            extraction_method=str(extraction.get("method") or "unknown"),
            pages_or_slides=extraction.get("pages_or_slides") or [],
            filter3_criteria=filter3_criteria,
            api_key=api_key,
            model=model,
        )

    status = "completed" if analysis.get("status") == "completed" else "completed_with_analysis_error"
    return {
        "document_id": document_id,
        "file_sha256": file_sha256,
        "filename": filename,
        "status": status,
        "extraction": extraction,
        "parser": parser_result,
        "deepseek_analysis": analysis,
        "viewer_conversion": conversion,
        "processed_at": utc_now(),
    }
